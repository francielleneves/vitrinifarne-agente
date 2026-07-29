"""Leitores de documentos da Vitrinifarne.

Uma função por formato de arquivo. Todas seguem o mesmo contrato:
recebem o caminho de um arquivo e devolvem uma string com o texto extraído.

Planilha, CSV e JSON são convertidos em frases (com o nome da coluna junto do
valor) para que a busca por similaridade consiga encontrá-los.

Uso:
    from leitores import extrair
    texto = extrair("docs/politica-de-privacidade.pdf", "pdf")
"""
import json
from pathlib import Path

import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


def ler_pdf(caminho):
    leitor = PdfReader(caminho)
    partes = []
    for numero, pagina in enumerate(leitor.pages, start=1):
        texto = (pagina.extract_text() or "").strip()
        if texto:
            partes.append(f"[Página {numero}]\n{texto}")
    return "\n\n".join(partes)


def ler_docx(caminho):
    doc = Document(caminho)
    partes = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for i, tabela in enumerate(doc.tables, start=1):
        linhas = [" | ".join(c.text.strip() for c in linha.cells) for linha in tabela.rows]
        partes.append(f"[Tabela {i}]\n" + "\n".join(linhas))
    return "\n\n".join(partes)


def ler_pptx(caminho):
    apresentacao = Presentation(caminho)
    partes = []
    for numero, slide in enumerate(apresentacao.slides, start=1):
        textos = []
        for forma in slide.shapes:
            if forma.has_text_frame and forma.text_frame.text.strip():
                textos.append(forma.text_frame.text.strip())
        if slide.has_notes_slide:
            nota = slide.notes_slide.notes_text_frame.text.strip()
            if nota:
                textos.append(f"Notas do apresentador: {nota}")
        if textos:
            partes.append(f"[Slide {numero}]\n" + "\n".join(textos))
    return "\n\n".join(partes)


def ler_xlsx(caminho):
    planilha = load_workbook(caminho, data_only=True)
    partes = []
    for aba in planilha.worksheets:
        linhas = list(aba.iter_rows(values_only=True))
        if not linhas:
            continue
        indice_cabecalho = None
        for i, linha in enumerate(linhas):
            preenchidas = sum(1 for v in linha if v not in (None, ""))
            if preenchidas >= 3:
                indice_cabecalho = i
                break
        if indice_cabecalho is None:
            continue
        cabecalho = [str(v).strip() if v is not None else "" for v in linhas[indice_cabecalho]]
        frases = []
        for linha in linhas[indice_cabecalho + 1:]:
            if all(v in (None, "") for v in linha):
                continue
            pares = [f"{cab}: {val}" for cab, val in zip(cabecalho, linha)
                     if cab and val not in (None, "")]
            if pares:
                frases.append("; ".join(pares))
        partes.append(f"[Aba: {aba.title}]\n" + "\n".join(frases))
    return "\n\n".join(partes)


def ler_csv(caminho):
    tabela = pd.read_csv(caminho)
    frases = []
    for _, linha in tabela.iterrows():
        pares = [f"{col}: {linha[col]}" for col in tabela.columns if pd.notna(linha[col])]
        frases.append("; ".join(pares))
    cabecalho = f"Tabela com {len(tabela)} registros e colunas: {', '.join(tabela.columns)}."
    return cabecalho + "\n" + "\n".join(frases)


def _achatar_json(dado, prefixo=""):
    linhas = []
    if isinstance(dado, dict):
        for chave, valor in dado.items():
            novo = f"{prefixo} > {chave}" if prefixo else str(chave)
            linhas.extend(_achatar_json(valor, novo))
    elif isinstance(dado, list):
        for i, item in enumerate(dado, start=1):
            novo = f"{prefixo} [{i}]" if prefixo else f"[{i}]"
            linhas.extend(_achatar_json(item, novo))
    else:
        linhas.append(f"{prefixo}: {dado}")
    return linhas


def ler_json(caminho):
    with open(caminho, encoding="utf-8") as arquivo:
        dado = json.load(arquivo)
    return "\n".join(_achatar_json(dado))


def ler_md(caminho):
    return Path(caminho).read_text(encoding="utf-8")


def ler_html(caminho):
    sopa = BeautifulSoup(Path(caminho).read_text(encoding="utf-8"), "html.parser")
    for tabela in sopa.find_all("table"):
        linhas_texto = []
        cabecalho = [th.get_text(strip=True) for th in tabela.find_all("th")]
        for tr in tabela.find_all("tr"):
            celulas = [td.get_text(strip=True) for td in tr.find_all("td")]
            if not celulas:
                continue
            if cabecalho and len(cabecalho) == len(celulas):
                linhas_texto.append("; ".join(f"{c}: {v}" for c, v in zip(cabecalho, celulas)))
            else:
                linhas_texto.append(" | ".join(celulas))
        tabela.replace_with("\n" + "\n".join(linhas_texto) + "\n")
    texto = sopa.get_text(separator="\n")
    return "\n".join(l.strip() for l in texto.splitlines() if l.strip())


LEITORES = {
    "pdf": ler_pdf, "docx": ler_docx, "pptx": ler_pptx, "xlsx": ler_xlsx,
    "csv": ler_csv, "json": ler_json, "markdown": ler_md, "html": ler_html,
}


def extrair(caminho, formato):
    leitor = LEITORES.get(formato)
    if leitor is None:
        raise ValueError(f"Formato sem leitor: {formato}")
    return leitor(caminho)
